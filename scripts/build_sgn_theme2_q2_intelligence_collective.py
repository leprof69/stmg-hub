# -*- coding: utf-8 -*-
"""Genere presentation SGN 1re STMG - Theme 2 Q2 intelligence collective."""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

DARK = RGBColor(0x16, 0x18, 0x28)
DARK_DEEP = RGBColor(0x0D, 0x0F, 0x18)
WHITE = RGBColor(0xF3, 0xF4, 0xF8)
MUTED = RGBColor(0xB0, 0xB8, 0xC8)
QUOTE_GOLD = RGBColor(0xFF, 0xD7, 0x4A)
BADGE_BG = RGBColor(0x2A, 0x2D, 0x3E)

N1 = RGBColor(0x29, 0xB6, 0xF6)
N2 = RGBColor(0x66, 0xBB, 0x6A)
N3 = RGBColor(0xFF, 0x98, 0x00)
N4 = RGBColor(0xAB, 0x47, 0xBC)

# Labels text only ASCII (dash for em-dash to avoid encoding issues in source)
BADGE_A = "A - EXERCICE"
BADGE_B = "B - CORRECTION"
BADGE_C = "C - CO-CONSTRUCTION"
BADGE_D = "D - TRACE ECRITE"


def _bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_badge(slide, text: str, accent: RGBColor | None = None, top: float = 0.12) -> None:
    w, h = Inches(1.42), Inches(0.42)
    left = Inches(10 - 1.42 - 0.25)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(top), w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = BADGE_BG
    if accent:
        shp.line.color.rgb = accent
        shp.line.width = Pt(2)
    else:
        shp.line.fill.background()
    tf = shp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = accent or QUOTE_GOLD
    p.font.name = "Calibri"


def _add_quote(slide, line: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.45), Inches(6.82), Inches(9.1), Inches(0.56))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"\u00ab {line} \u00bb"
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = QUOTE_GOLD
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER


def _title_box(slide, title: str, color: RGBColor = WHITE, size: int = 26) -> None:
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.08), Inches(8.0), Inches(0.92))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Calibri Light"


def _para_markdown(
    paragraph,
    text: str,
    *,
    size: Pt,
    color: RGBColor,
    default_bold: bool = False,
) -> None:
    paragraph.clear()
    for i, part in enumerate(text.split("**")):
        if not part:
            continue
        run = paragraph.add_run()
        run.text = part
        run.font.name = "Calibri"
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = default_bold or (i % 2 == 1)


def _body_box(
    slide,
    lines: list[str],
    top: float = 1.12,
    color: RGBColor = MUTED,
    size: int = 18,
) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9.0), Inches(5.48))
    tf = box.text_frame
    tf.word_wrap = True
    pt = Pt(size)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _para_markdown(p, line, size=pt, color=color)
        p.space_after = Pt(10)


def _table_slide(
    slide,
    rows: list[tuple[str, str]],
    top: float = 1.05,
    accent: RGBColor = N1,
) -> None:
    nrows = len(rows)
    table = slide.shapes.add_table(
        nrows,
        2,
        Inches(0.45),
        Inches(top),
        Inches(9.1),
        Inches(min(5.4, 0.47 * nrows + 1.55)),
    ).table
    for r, (a, b) in enumerate(rows):
        ca = table.cell(r, 0)
        cb = table.cell(r, 1)
        ha = ca.text_frame.paragraphs[0]
        hb = cb.text_frame.paragraphs[0]
        fz = Pt(13 if r > 0 else 13)
        col = DARK_DEEP if r == 0 else WHITE
        _para_markdown(ha, a, size=fz, color=col, default_bold=(r == 0))
        _para_markdown(hb, b, size=fz, color=col, default_bold=(r == 0))
        table.rows[r].height = Pt(46 if r > 0 else 32)
        if r == 0:
            ca.fill.solid()
            cb.fill.solid()
            ca.fill.fore_color.rgb = accent
            cb.fill.fore_color.rgb = accent


def _rect_note(slide, lines: list[str], top: float, accent: RGBColor) -> None:
    h = Inches(min(2.5, 0.3 + 0.2 * len(lines)))
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.48), Inches(top), Inches(9.05), h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x22, 0x25, 0x38)
    box.line.color.rgb = accent
    box.line.width = Pt(2)
    tf = box.text_frame
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)
    pt = Pt(15)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _para_markdown(p, line, size=pt, color=WHITE, default_bold=(i == 0))


def _schema_recap(slide) -> None:
    txt = (
        "[Partage des infos au sein ORGA]\n"
        "|-- E-com (bleu) : SI trace + mails / ENT / MI / visio\n"
        "|-- Collaboration (vert) : cloud temps reel + outils coproduit\n"
        "|-- Communautes / RSE (orange) : memoire informelle talents\n"
        "|-- IA (violet) : automatise repetitifs + amplify humains\n"
        "=> **Intelligence collective** = capitaliser savoir distribue\n"
        "             + coherence des decisions ensemble"
    )
    box = slide.shapes.add_textbox(Inches(0.42), Inches(1.0), Inches(9.16), Inches(5.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = txt
    p.font.name = "Consolas"
    p.font.size = Pt(15)
    p.font.color.rgb = WHITE


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_slide():
        s = prs.slides.add_slide(blank)
        _bg(s, DARK)
        return s

    q_gest = (
        "Comment le partage de l'information contribue-t-il \u00e0 l'\u00e9mergence d'une intelligence collective ?"
    )

    # 1
    s = add_slide()
    _add_badge(s, "OUVERTURE", QUOTE_GOLD)
    _title_box(
        s,
        "SGN 1\u00e8re STMG - Theme 2\u00a0: episode special cohesion MAX (\u266b)",
        N1,
        22,
    )
    _body_box(
        s,
        [
            "Question :\n" + q_gest,
            "4 blocs (programme officiel Theme 2) + cas manga / DA + exemples reel entreprise.",
            "Charte slide : fond sombre, couleurs vives, citations anim\u00e9s sous chaque planche.",
        ],
        top=1.38,
        size=16,
        color=WHITE,
    )
    _add_quote(s, "Tu bloques sur le boss ? Branche tes alli\u00e9s."

                 " Partage tes infos.")

    # 2
    s = add_slide()
    _add_badge(s, "ACCROCHE", QUOTE_GOLD)
    _title_box(s, "Team 7, Z Warriors, Totally Spies : m\u00eame recette ?", QUOTE_GOLD, 24)
    _body_box(
        s,
        [
            ("Intelligence collective = info partag\u00e9e + confiance "
             "+ coordination vers objectif commun (plus vite qu'en solo isol\u00e9)."),
            "Le numero massivement interconnecte cet \u00e9change d'infos tout en rajoutant des risques SI/cybers\u00e9curit\u00e9/RGPD.",
            "Hypoth\u00e8se cours : quatre familles de dispositifs orga donnent corps \u00e0 cette question.",
            "On connecte ensuite chaque famille \u00e0 tes dessins anim\u00e9s pr\u00e9f\u00e9r\u00e9s.",
        ],
        top=1.25,
        size=16,
        color=WHITE,
    )
    _add_quote(s, "Personne ne sauve Paris seul alors que le super-vilain multit\u00e2che.")

    # N1 slides 3-6
    s = add_slide()
    _add_badge(s, BADGE_A, N1)
    _title_box(s, "\u00c0 toi de jouer ! - Capsule Corp (Dragon Ball)", N1)
    _body_box(
        s,
        [
            ("Brief :\n"
             "Siege hyper-tech coordonne usines + franchises mondiales :\n"
             "il faut des flux synchro avant le prochain duel industriel."),
            "1 Une donn\u00e9e brute sert quand elle devient informations puis... ?",
            "2 Donnes 3 canaux numero pour relier vite les centres partout sur la plan\u00e8te.",
            "3 Diff\u00e9rences internes (equipe siege) VS externes (boutiques licences) dans les \u00e9changes.",
            "4 Risque x 2 avec messageries ENT mal gouvern\u00e9es (confidentialite, surcharge,...).",
        ],
        top=0.92,
        size=14,
        color=MUTED,
    )
    _add_quote(s, "Capte tous les signaux avant d'attaquer.")

    s = add_slide()
    _add_badge(s, BADGE_B, N1)
    _title_box(s, "Correction - E-communication et partage d'infos", WHITE)
    _body_box(
        s,
        [
            "**E-communication** : \u00e9changes d'infos par outils numero internes / partenaires (mail, MI, visio,...).",
            "Triptyque programme : **donn\u00e9e** \u2192 **information** \u2192 **connaissance** partagee (voir QTheme2-1).",
            "**Communication interne** vs **externe** ; **SI** classe, diffuse, trace ; **PGI**/ENT soutiennent.",
            "Atouts : **rapidite**, accessibilite, **cout** communication,\n **tracabilite**.",
            ("Risques : **surcharge** info, pertes **confidentialite** / RGPD,"
             " **dependance** techno, incidents **cybersecurite**."),
        ],
        top=0.94,
        size=13,
        color=WHITE,
    )
    _add_quote(s, "Le combat se gagne aussi avec des lignes propres et trac\u00e9es.")

    s = add_slide()
    _add_badge(s, BADGE_C, N1)
    _title_box(s, "Construisons ensemble - E-communication", N1)
    _body_box(
        s,
        [
            "D'apres le cas : le numero sert surtout \u00e0 transformer... en... ?",
            "Ce que j'ai compris : **SI** = moteur de circulation car...",
            "Interne / externe : la fronti\u00e8re utile \u00e0 retenir c'est...",
            "Tracer les \u00e9changes, c'est aussi prot\u00e9ger... (finis).",
        ],
        top=1.2,
        size=16,
    )
    _add_quote(s, "Briefing clair au bon timing = demi victoire avant l'opening.")

    s = add_slide()
    _add_badge(s, BADGE_D, N1)
    _title_box(s, "Trace ecrite - E-communication", N1)
    _table_slide(
        s,
        [
            (
                "Definition officielle",
                "**E-communication** regroupe tous \u00e9changes infos via outils numero **internes** ou **externes**\n"
                "(*mail, messagerie instantan\u00e9e, **visio**, ENT, intranet, extranet*). Transformation **donnee \u2192 info "
                "\u2192 savoir mutualise**, port\u00e9 par le **systeme d'information** (**PGI**,\n"
                "**ENT**, plateformes collaboratives,...). Programme SDGN 1re Theme 2 applications numeriques organisations.",
            ),
            (
                "Exemple reel | DA",
                "Grands groupes type **aero** (ENT + visio + PGI) vs **Capsule Corp** briefings scouter plan\u00e8te.",
            ),
            (
                "A retenir",
                "**Rapidite + accessibilite + trace** ; risques **surcharge / confidentialite / cyber**.",
            ),
        ],
        top=0.88,
        accent=N1,
    )
    _rect_note(
        s,
        [
            "A retenir (interro)",
            "\u2022 Formes citees \u2022 Flux internes/externes \u2022 Role **SI**/PGI/**ENT**",
        ],
        top=4.95,
        accent=N1,
    )
    _add_quote(s, "Capsule envoyee hors timing = info perdue avant Zarbon.")

    # N2 slides 7-10
    s = add_slide()
    _add_badge(s, BADGE_A, N2)
    _title_box(s, "\u00c0 toi de jouer ! - Thousand Sunny (One Piece)", N2)
    _body_box(
        s,
        [
            ("Equipage eclate : certains infiltr\u00e9s, autres en mer ;"
             " tous doivent co-\u00e9crire le plan contre la Marine contemporaine."),
            "1 Qui recoit passivement infos VS qui **coproduit** version finale document ? Explique vite.",
            "2 Donne trois familles d'outils **temps quasi reel** (docs, tableau, salons vocaux,...).",
            "3 Coordination numero : cite 2 b\u00e9n\u00e9fices + 2 limites humaines (isolement,...).",
            "4 Pourquoi le **cloud** change lieu / moment ou on travaille ?",
        ],
        top=1.0,
        size=14,
    )
    _add_quote(s, "Une feuille de route vivante vaut mille monologues \u00e0 la passerelle.")

    s = add_slide()
    _add_badge(s, BADGE_B, N2)
    _title_box(s, "Correction - Collaboration et outils collaboratifs", WHITE)
    _body_box(
        s,
        [
            "**Collaboration** = objectif commun + **coproduction** info / decision / livrable (**Google Docs**, **Office365**).",
            "Outils agendas partag\u00e9s, wikis, **Slack**/Teams, Kanban (**Trello**/Notion) \u2260 simple mail passif.",
            "**Cloud computing** = stockage synchro depuis tout terminal + droits granularit\u00e9.",
            "**Teletravail** / mobilit\u00e9 gagnent flexibilite temps ; risque lien social amoindri coordination humaine delicate.",
            "Cap : fracture numerique ou securite hebergeur si cloud mal param\u00e9tr\u00e9 / pas de RGPD.",
        ],
        top=0.94,
        size=13,
        color=WHITE,
    )
    _add_quote(s, "Alliance Pirate : la tactique commune sur \u00e9cran fait tenir les Nakama.")

    s = add_slide()
    _add_badge(s, BADGE_C, N2)
    _title_box(s, "Construisons ensemble - Collaboration", N2)
    _body_box(
        s,
        [
            "Collaboration VS simple briefing : moi je retiens une difference cle = ...",
            "Deux usages cloud personnels que vous connaissez tous : ... + ...",
            "Si personne ne se voit synchro trois semaines, je redoute ... pourquoi ?",
            "Phrase souvenir : plusieurs cerveaux + **meme tableau** = intelligence collective mieux distribu\u00e9e car...",
        ],
        top=1.08,
        size=15,
    )
    _add_quote(s, "Pas d'Alliance sans planning partag\u00e9 vivant sous les yeux de tout le monde.")

    s = add_slide()
    _add_badge(s, BADGE_D, N2)
    _title_box(s, "Trace ecrite - Collaboration", N2)
    _table_slide(
        s,
        [
            (
                "Definition officielle",
                "**Collaboration**\n Travail conjoint **internes/externes**\n Vers **but commun**, facilite numero "
                "**coproduit** informations / decisions / livrables (**agendas**\n docs cloud, Wikis,...). Travail distant / "
                "**mobilit\u00e9** + infra **cloud** (programme officiel mentions).",
            ),
            (
                "Exemple reel | DA",
                "Projets Renault + suites **Office365**\nVers ** Mugiwara ** roadmaps vivantes depuis telephones navire.",
            ),
            (
                "A retenir",
                "Cloud = **infra collab**\n Mais humains > simple stack Slack.\nLimiter fractures / risques tiers.",
            ),
        ],
        top=0.9,
        accent=N2,
    )
    _rect_note(
        s,
        ["A retenir", "\u2022 Collaborer vs informer \u2022 Nuage indispensable \u2022 Risque isolement managerial"],
        top=5.0,
        accent=N2,
    )
    _add_quote(s, "La victoire se code \u00e0 plusieurs bras sur meme console nuage.")

    # N3 11-14
    s = add_slide()
    _add_badge(s, BADGE_A, N3)
    _title_box(s, "\u00c0 toi de jouer ! - Ligue Pokemon + WOOHP (Totally Spies)", N3)
    _body_box(
        s,
        [
            "Canal ferm\u00e9 : dresseurs + staff echangent combos, raids, anomalies biome instantan\u00e9.",
            ("1 Trouve une etiquette conceptuelle (\u2260 titre du cours encore) pour "
             "petit groupe continu en ligne meme passion."),
            ("2 Pourquoi ce flux n'est pas simplement votre Insta personnel grand public,"
             "\nmais un espace cloison\u00e9 type QG ?"),
            "3 Liste deux bienfaits ORGA lorsque feedback terrain remonte spontan\u00e9.",
            ("4 Risque management si salons informels proliferent sans moderation"
             "(temps, rumeurs, productivite, e-reputation interne)."),
        ],
        top=0.94,
        size=13,
    )
    _add_quote(s, "Chaque badge raconte ses fails : niveau groupe augmente vite sans ego bloquant.")

    s = add_slide()
    _add_badge(s, BADGE_B, N3)
    _title_box(s, "Correction - Communautes en ligne et RSE", WHITE)
    _body_box(
        s,
        [
            "**Communaute en ligne** : gens connect\u00e9s recurrentment autour int\u00e9r\u00eat / projet commun.",
            "**RSE** (**Teams/Yammer/Slack/Workplace**) = reseau social **entreprise**\nSavoir tacite diffuse et visible.",
            "Distinction FB/Insta/Linked **public**\nOu mixte perso : finalites / visibilit\u00e9s differentes managers.",
            "**Communaute de pratiques** : partage tacite exp\u00e9riences, bonnes pratiques spontan\u00e9.",
            "**Lien intelligence collective** :\n Agr\u00e8ge savoirs individuels en memo collective organisationnelle participative.",
        ],
        top=0.94,
        size=13,
        color=WHITE,
    )
    _add_quote(s, "Pokemon Center + fil mission WOOHP = savoir tacite vite capitalise.")

    s = add_slide()
    _add_badge(s, BADGE_C, N3)
    _title_box(s, "Construisons ensemble - Communautes / RSE", N3)
    _body_box(
        s,
        [
            ("Quand des experts postent sans hierarchie pyramidale,"
             " moi j'appele cet espace communautaire car..."),
            "La plateforme permet de faire emerger quel type de memoire comparativement au mail?",
            "**Pro**\nversus **grand public**\nOu je trace la ligne nette :\n____________",
            "Connaissance individuelle -\u003e ??? quand plusieurs likes / commentaires structurent reponse.",
        ],
        top=1.02,
        size=15,
    )
    _add_quote(s, "Totally Spies : debriefs internes > stories publiques pour sauver la mission.")

    s = add_slide()
    _add_badge(s, BADGE_D, N3)
    _title_box(s, "Trace ecrite - Communautes en ligne et RSE", N3)
    _table_slide(
        s,
        [
            (
                "Definition officielle",
                "**Communaute en ligne**\n Individus interagissant souvent numero autour meme pratiques / ambitions.\n"
                "**Reseaux sociaux d'entreprise (RSE)** : plate-formes internes echanges informels,"
                "**bonnes pratiques**, cohesion, innovation communautaire. Programme officiel cites.",
            ),
            (
                "Exemple reel | DA",
                "**Yammer** + salons **Teams** vs **clubs Pokemon**\n+\nFil **WOOHP** :\ncloisonnement pro.",
            ),
            (
                "A retenir",
                "**Bonnes pratiques** + cohesion + memo ; veiller dispersion / **rumeurs** / temps perdu.",
            ),
        ],
        top=0.92,
        accent=N3,
    )
    _rect_note(
        s,
        ["A retenir", "\u2022 **RSE** internes \u2022 Communaute pratiques \u2022 Lien memo collective IQ orga"],
        top=5.0,
        accent=N3,
    )
    _add_quote(s, "Des dizaines feedback terrain = strategist brain upgrade pour le QG.")

    # N4 slides 15-18
    s = add_slide()
    _add_badge(s, BADGE_A, N4)
    _title_box(s, "\u00c0 toi de jouer ! - Tony Stark (MCU) et Dr Stone", N4)
    _body_box(
        s,
        [
            ("Labo : robot analyse historique commandes, predit pannes,"
             " propose planning entretiens ; modelise cas clients."),
            "1 Regle SI/ALORS figee VS systeme qui apprend puis re-ordonne tout seul ?",
            "2 Cite deux taches ORGA deleguables bots + deux ou la relation humaine reste roi.",
            "3 Productivite gagnee -\u003e quels impacts sur metiers repetitifs / nouveaux metiers data ?",
            "4 Risque ethic / RGPD concret avec ce labo automatise ?",
        ],
        top=0.94,
        size=14,
    )
    _add_quote(s, "Friday, aide-moi a synthetiser, mais nous validons tous ensemble avant agir.")

    s = add_slide()
    _add_badge(s, BADGE_B, N4)
    _title_box(s, "Correction - IA et automatisation des taches organisations", WHITE)
    _body_box(
        s,
        [
            "**IA** mime apprentissage / raisonnement / aide-decision vs **automate** qui rejoue script fixe.",
            "Exemples : chatbots SAV ; reco e-commerce ; **analyse predictive** RH/finance ; vision / controle qualite.",
            "Automatisable : repetitions / peu valeur ajoutee ; delicate : empathie,\n**\u2013**\ncreativite,\ndecisions ethiques complexes.",
            "Travail : productivites + metamorphoses emplois (data scientist, prompt engineer,...).",
            "Ethiques : **biais**\n**\u2013**\n**transparence**\n**\u2013**\nRGPD\n**\u2013**\nculpabilite erreur algo.",
        ],
        top=0.94,
        size=13,
        color=WHITE,
    )
    _add_quote(s, "Senku : la science acc\u00e8l\u00e8re, nos choix sociaux restent humains.")

    s = add_slide()
    _add_badge(s, BADGE_C, N4)
    _title_box(s, "Construisons ensemble - IA et automatisation", N4)
    _body_box(
        s,
        [
            "Une **IA** n'est pas automatiquement intelligence collective :\n moi je vois \u00e7a comme outil servant \u00e0...",
            "Si un algorithme biaise le recrutement : qui dois-tu rendre accountable c\u00f4te organisation ?",
            "Metier nouveau lie IA que tu cites en veille :\n____________",
            "Donne trois questions ethiques indispensables avant de deployer bots clients.",
        ],
        top=1.05,
        size=15,
    )
    _add_quote(s, "La machine compile ; l'equipe decide du sens commun.")

    s = add_slide()
    _add_badge(s, BADGE_D, N4)
    _title_box(s, "Trace ecrite - IA et automatisation organisationnelles", N4)
    _table_slide(
        s,
        [
            (
                "Definition officielle",
                "**Intelligence artificielle** :\n syst\u00e8me mimant cognition (apprend, raisonne, conseille)."
                "**Automatisation** orga :\ndeleguer repetitions /\nanalyser volumetries pour gagner temps."
                "Programme :\n usages numeriques -\u003epilotage -\u003e limites RGPD/ethique.",
            ),
            (
                "Exemple reel | DA",
                "Chatbots + reco industriels -\u003evs **Friday / labos Senku** assistent brainstorming humain.",
            ),
            (
                "A retenir",
                "**Levier**\n intelligence collective :\namplify humains \u2260 remplacer **co-construction**\n culturelle.",
            ),
        ],
        top=0.9,
        accent=N4,
    )
    _rect_note(
        s,
        [
            "A retenir",
            "\u2022 IA adapte VS automate fixe \u2022 Metiers remixes \u2022 RGPD + biais + responsabilites",
        ],
        top=5.0,
        accent=N4,
    )
    _add_quote(s, "Technologie de ouf :\n**\u2666**\nmais aucun succes durable sans pacte moral equipe.")

    # 19 - Conclusion carte mentale
    s = add_slide()
    _add_badge(s, "FIN EPISODE", QUOTE_GOLD)
    _title_box(s, "Carte finale : quatre leviers -\u003e meme question gestion ?", QUOTE_GOLD, 22)
    _schema_recap(s)
    _add_quote(s, "Generique :\n**\u2666**\nrevois ton cahier puis enchaine revisions Theme 3 tranquille.")

    # 20 - Bilan
    s = add_slide()
    _add_badge(s, "DEBAT FINAL", QUOTE_GOLD)
    _title_box(s, "Sans lien humain :\n**\u2666**\npeut-on maintenir une intelligence collective forte ?", WHITE, 20)
    _body_box(
        s,
        [
            "**Motion** : startup 100% distante + bots = m\u00eame esprit commun qu'un resto physique ? Pour / contre vite.",
            "Relie tes arguments aux quatre couches : notions bleues, vertes, oranges, puis violettes du cours.",
            "Devoir (\u2264 cinq lignes) : synth\u00e8se finale \u00e0 partir des traces \u00e9crites + argument perso mesur\u00e9.",
        ],
        top=1.55,
        size=15,
        color=WHITE,
    )
    _add_quote(s, "A suivre : humains -\u003enumerique :\n meilleur duo shonen quand ils se parlent encore.")

    return prs


def main() -> None:
    import os

    repo = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    exports = os.path.join(repo, "exports")
    os.makedirs(exports, exist_ok=True)
    out = os.path.join(exports, "SGN_1re_Theme2_Q2_Intelligence_collective_DA.pptx")
    build().save(out)
    print("OK ->", out)


if __name__ == "__main__":
    main()
